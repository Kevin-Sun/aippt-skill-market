#!/usr/bin/env python3
"""解析 .pptx 文件提取版式结构/配色/字体（P0 修复：layouts 解析 + palettes 提取）"""
import sys, json, os
from pptx import Presentation
from pptx.util import Emu

def parse(pptx_path):
    result = {"file": pptx_path, "parsed": False, "slides": [], "palettes": [], "fonts": [], "error": None}
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        result["error"] = f"open_fail: {e}"
        return result

    colors = set()
    fonts = set()
    layout_names = set()

    for idx, slide in enumerate(prs.slides):
        slide_info = {"index": idx, "layout": slide.slide_layout.name, "shapes": []}
        layout_names.add(slide.slide_layout.name)
        for shape in slide.shapes:
            shape_info = {"type": str(shape.shape_type), "name": shape.name}
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            fonts.add(run.font.name)
                        try:
                            if run.font.color and run.font.color.rgb:
                                colors.add(str(run.font.color.rgb))
                        except Exception:
                            pass
            if hasattr(shape, "fill"):
                try:
                    if shape.fill.type is not None and shape.fill.fore_color.rgb:
                        colors.add(str(shape.fill.fore_color.rgb))
                except Exception:
                    pass
            slide_info["shapes"].append(shape_info)
        result["slides"].append(slide_info)

    result["palettes"] = sorted(colors)[:20]
    result["fonts"] = sorted(fonts)[:10]
    result["layoutNames"] = sorted(layout_names)
    result["slideCount"] = len(prs.slides)
    result["parsed"] = True
    return result

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"error": "usage: parse_pptx.py <pptx_path> [output_json] [source] [license] [category]"}))
        sys.exit(1)
    pptx_path = sys.argv[1]
    out_path = sys.argv[2] if len(sys.argv) > 2 else None
    source = sys.argv[3] if len(sys.argv) > 3 else ""
    license_ = sys.argv[4] if len(sys.argv) > 4 else "unknown"
    category = sys.argv[5] if len(sys.argv) > 5 else ""
    result = parse(pptx_path)
    result["source"] = source
    result["license"] = license_
    result["category"] = category
    if out_path:
        with open(out_path, "w") as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
    else:
        print(json.dumps(result, ensure_ascii=False))
